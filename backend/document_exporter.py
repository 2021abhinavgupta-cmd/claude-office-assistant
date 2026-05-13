"""
document_exporter.py — Server-side document export
Generates production-quality DOCX, PDF, and PPTX from Markdown.
Uses python-docx for Word, WeasyPrint for PDF, python-pptx for PowerPoint.
"""

import io
import logging
import os
import re
import subprocess
import sys
import html as html_module
from typing import List, Optional, Tuple

logger = logging.getLogger(__name__)

_pptx_pip_attempted = False


def _pptx_install_help() -> str:
    return (
        "PowerPoint export needs the python-pptx package. "
        "Install it with: pip install python-pptx "
        "(from your project folder: pip install -r backend/requirements.txt). "
        "Use the same Python as your server (e.g. venv/bin/python)."
    )


def _ensure_python_pptx_installed() -> None:
    """
    Ensure `pptx` can be imported. If missing, try `pip install python-pptx` using **this**
    interpreter (matches gunicorn / `python app.py`). Set OFFICE_ASSISTANT_DISABLE_AUTO_PIP_PPTX=1
    to skip auto-install (locked-down hosts).
    """
    global _pptx_pip_attempted
    try:
        import pptx  # noqa: F401
        return
    except ImportError:
        pass

    if os.getenv("OFFICE_ASSISTANT_DISABLE_AUTO_PIP_PPTX", "").strip().lower() in (
        "1",
        "true",
        "yes",
        "on",
    ):
        raise RuntimeError(_pptx_install_help())

    if _pptx_pip_attempted:
        raise RuntimeError(_pptx_install_help())

    _pptx_pip_attempted = True
    logger.warning(
        "python-pptx not importable — running pip install into %s",
        sys.executable,
    )
    try:
        proc = subprocess.run(
            [
                sys.executable,
                "-m",
                "pip",
                "install",
                "python-pptx>=0.6.23",
            ],
            check=False,
            timeout=180,
            capture_output=True,
            text=True,
        )
        if proc.returncode != 0:
            logger.error(
                "pip install python-pptx failed (exit %s): %s",
                proc.returncode,
                (proc.stderr or proc.stdout or "")[:2000],
            )
            raise RuntimeError(_pptx_install_help())
    except subprocess.TimeoutExpired as e:
        logger.error("pip install python-pptx timed out")
        raise RuntimeError(_pptx_install_help()) from e
    except OSError as e:
        logger.error("could not run pip: %s", e)
        raise RuntimeError(_pptx_install_help()) from e

    import importlib

    importlib.invalidate_caches()
    try:
        import pptx  # noqa: F401
    except ImportError as e:
        logger.error("pptx still not importable after pip install")
        raise RuntimeError(_pptx_install_help()) from e

# ─── DOCX Export ─────────────────────────────────────────────────────────────

def export_docx(markdown_text: str, title: str = "Document") -> io.BytesIO:
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # ── Page margins ──────────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin    = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin   = Inches(1.15)
        section.right_margin  = Inches(1.15)

    # ── Styles ────────────────────────────────────────────────────────────────
    styles = doc.styles

    # Normal text
    normal = styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(6)

    # Heading styles
    for lvl, (size, bold) in enumerate([(20, True), (16, True), (13, True), (12, True)], start=1):
        try:
            hs = styles[f"Heading {lvl}"]
            hs.font.name  = "Calibri"
            hs.font.size  = Pt(size)
            hs.font.bold  = bold
            hs.font.color.rgb = RGBColor(0x1a, 0x56, 0xdb)
            hs.paragraph_format.space_before = Pt(14)
            hs.paragraph_format.space_after  = Pt(4)
        except Exception:
            pass

    # ── Helpers ───────────────────────────────────────────────────────────────
    def add_formatted_run(paragraph, text: str):
        """Parse inline **bold**, *italic*, `code` and add styled runs."""
        pattern = re.compile(r"(\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`)")
        last = 0
        for m in pattern.finditer(text):
            # Plain text before match
            if m.start() > last:
                paragraph.add_run(text[last:m.start()])
            full = m.group(0)
            if full.startswith("**"):
                r = paragraph.add_run(m.group(2))
                r.bold = True
            elif full.startswith("*"):
                r = paragraph.add_run(m.group(3))
                r.italic = True
            elif full.startswith("`"):
                r = paragraph.add_run(m.group(4))
                r.font.name = "Courier New"
                r.font.size = Pt(10)
            last = m.end()
        if last < len(text):
            paragraph.add_run(text[last:])

    def add_table_row(table, cells: list, is_header: bool = False):
        from docx.shared import RGBColor
        row = table.add_row()
        for i, cell_text in enumerate(cells):
            if i >= len(row.cells):
                break
            cell = row.cells[i]
            p = cell.paragraphs[0]
            run = p.add_run(cell_text.strip())
            if is_header:
                run.bold = True
                shading = OxmlElement("w:shd")
                shading.set(qn("w:val"), "clear")
                shading.set(qn("w:color"), "auto")
                shading.set(qn("w:fill"), "D9E1F2")
                cell._tc.get_or_add_tcPr().append(shading)

    # ── Parse & Render Markdown ───────────────────────────────────────────────
    lines   = markdown_text.splitlines()
    i       = 0
    in_code = False
    code_buf= []

    while i < len(lines):
        line = lines[i]

        # Code block
        if line.strip().startswith("```"):
            if not in_code:
                in_code  = True
                code_buf = []
                i += 1
                continue
            else:
                in_code = False
                p = doc.add_paragraph()
                p.style = "No Spacing"
                run = p.add_run("\n".join(code_buf))
                run.font.name = "Courier New"
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(0x16, 0x3a, 0x5f)
                shd = OxmlElement("w:shd")
                shd.set(qn("w:val"),   "clear")
                shd.set(qn("w:color"), "auto")
                shd.set(qn("w:fill"),  "F0F4F8")
                p._p.get_or_add_pPr().append(shd)
                code_buf = []
                i += 1
                continue

        if in_code:
            code_buf.append(line)
            i += 1
            continue

        # Horizontal rule
        if re.match(r"^-{3,}$|^\*{3,}$|^_{3,}$", line.strip()):
            doc.add_paragraph("─" * 60, style="Normal")
            i += 1
            continue

        # Headings
        m_h = re.match(r"^(#{1,4})\s+(.*)", line)
        if m_h:
            lvl = len(m_h.group(1))
            p = doc.add_heading(m_h.group(2).strip(), level=lvl)
            i += 1
            continue

        # Table — detect by | at start
        if "|" in line and line.strip().startswith("|"):
            # Collect all table lines
            table_lines = []
            while i < len(lines) and "|" in lines[i] and lines[i].strip().startswith("|"):
                if not re.match(r"^\|[-| :]+\|$", lines[i].strip()):
                    table_lines.append(lines[i])
                i += 1
            if not table_lines:
                continue
            headers  = [c.strip() for c in table_lines[0].strip("|").split("|")]
            data_rows = [
                [c.strip() for c in row.strip("|").split("|")]
                for row in table_lines[1:]
            ]
            max_cols = max(len(headers), max((len(r) for r in data_rows), default=0))
            tbl = doc.add_table(rows=1, cols=max_cols)
            tbl.style = "Table Grid"
            add_table_row(tbl, headers, is_header=True)
            for row in data_rows:
                # Pad to max_cols
                row += [""] * (max_cols - len(row))
                add_table_row(tbl, row)
            doc.add_paragraph("")
            continue

        # Unordered list
        m_ul = re.match(r"^(\s*)[-*+]\s+(.*)", line)
        if m_ul:
            indent = len(m_ul.group(1)) // 2
            p = doc.add_paragraph(style="List Bullet" if indent == 0 else "List Bullet 2")
            add_formatted_run(p, m_ul.group(2))
            i += 1
            continue

        # Ordered list
        m_ol = re.match(r"^(\s*)\d+\.\s+(.*)", line)
        if m_ol:
            p = doc.add_paragraph(style="List Number")
            add_formatted_run(p, m_ol.group(2))
            i += 1
            continue

        # Blockquote
        m_bq = re.match(r"^>\s*(.*)", line)
        if m_bq:
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0.4)
            run = p.add_run(m_bq.group(1))
            run.italic = True
            run.font.color.rgb = RGBColor(0x60, 0x60, 0x60)
            i += 1
            continue

        # Empty line
        if not line.strip():
            i += 1
            continue

        # Normal paragraph
        p = doc.add_paragraph()
        add_formatted_run(p, line.strip())
        i += 1

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ─── PDF Export ──────────────────────────────────────────────────────────────


def export_pdf_reportlab(markdown_text: str, title: str = "Document") -> io.BytesIO:
    """
    Pure-Python PDF fallback when WeasyPrint is missing or fails (e.g. minimal Railway image).
    Renders markdown-ish content as readable paragraphs (no full HTML styling).
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
    from xml.sax.saxutils import escape

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        title=(title or "Export")[:120],
        leftMargin=54,
        rightMargin=54,
        topMargin=54,
        bottomMargin=54,
    )
    styles = getSampleStyleSheet()
    story: List[object] = []

    story.append(Paragraph(f"<b>{escape(title or 'Document')}</b>", styles["Title"]))
    story.append(Spacer(1, 16))

    raw = (markdown_text or "").strip()
    if not raw:
        story.append(Paragraph("<i>(empty)</i>", styles["Normal"]))
        doc.build(story)
        buf.seek(0)
        return buf

    blocks = re.split(r"\n\s*\n+", raw)
    for blk in blocks:
        for line in blk.splitlines():
            ln = line.strip()
            if not ln:
                continue
            if ln.startswith("```"):
                continue
            if re.match(r"^```\w*$", ln):
                continue
            # Heading lines → bold
            if re.match(r"^#{1,6}\s+", ln):
                inner = re.sub(r"^#{1,6}\s+", "", ln).strip()
                story.append(Paragraph(f"<b>{escape(inner)}</b>", styles["Heading2"]))
                continue
            # Bullet / numbered
            if re.match(r"^[-*+]\s+", ln):
                inner = re.sub(r"^[-*+]\s+", "", ln).strip()
                story.append(Paragraph(f"• {escape(inner)}", styles["Normal"]))
                continue
            if re.match(r"^\d+\.\s+", ln):
                inner = re.sub(r"^\d+\.\s+", "", ln).strip()
                story.append(Paragraph(escape(inner), styles["Normal"]))
                continue
            story.append(Paragraph(escape(ln).replace("\n", "<br/>"), styles["Normal"]))
        story.append(Spacer(1, 8))

    doc.build(story)
    buf.seek(0)
    return buf


def export_pdf(markdown_text: str, title: str = "Document") -> io.BytesIO:
    try:
        import markdown as md_lib
        html_body = md_lib.markdown(
            markdown_text,
            extensions=["tables", "fenced_code", "nl2br"]
        )
    except ImportError:
        # Fallback: simple HTML from text
        escaped = html_module.escape(markdown_text).replace("\n", "<br>")
        html_body = f"<p>{escaped}</p>"

    css = """
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700&family=JetBrains+Mono&display=swap');
    * { box-sizing: border-box; margin: 0; padding: 0; }
    body {
        font-family: 'Inter', 'Segoe UI', Arial, sans-serif;
        font-size: 11pt;
        line-height: 1.7;
        color: #1a1a2e;
        padding: 48px 60px;
        max-width: 900px;
        margin: 0 auto;
    }
    h1 { font-size: 22pt; font-weight: 700; color: #1a56db; margin: 24px 0 10px; }
    h2 { font-size: 16pt; font-weight: 700; color: #1a56db; margin: 20px 0 8px; border-bottom: 2px solid #e5eaf5; padding-bottom: 4px; }
    h3 { font-size: 13pt; font-weight: 600; color: #2d3a6b; margin: 16px 0 6px; }
    h4 { font-size: 11pt; font-weight: 600; color: #2d3a6b; margin: 12px 0 4px; }
    p  { margin-bottom: 10px; }
    ul, ol { margin: 8px 0 8px 24px; }
    li { margin-bottom: 4px; }
    strong { font-weight: 600; }
    em { font-style: italic; color: #444; }
    code {
        font-family: 'JetBrains Mono', 'Courier New', monospace;
        background: #f0f4f8;
        border: 1px solid #dde3ee;
        border-radius: 4px;
        padding: 2px 6px;
        font-size: 9.5pt;
    }
    pre {
        background: #f0f4f8;
        border: 1px solid #dde3ee;
        border-radius: 8px;
        padding: 16px;
        margin: 12px 0;
        overflow-x: auto;
    }
    pre code { background: none; border: none; padding: 0; }
    blockquote {
        border-left: 4px solid #1a56db;
        margin: 12px 0 12px 0;
        padding: 8px 16px;
        color: #555;
        font-style: italic;
        background: #f5f8ff;
    }
    table {
        border-collapse: collapse;
        width: 100%;
        margin: 14px 0;
        font-size: 10pt;
    }
    th {
        background: #1a56db;
        color: #fff;
        padding: 10px 12px;
        font-weight: 600;
        text-align: left;
    }
    td {
        border: 1px solid #dde3ee;
        padding: 8px 12px;
        vertical-align: top;
    }
    tr:nth-child(even) td { background: #f5f8ff; }
    hr { border: none; border-top: 2px solid #e5eaf5; margin: 20px 0; }
    .title-block {
        border-bottom: 3px solid #1a56db;
        margin-bottom: 24px;
        padding-bottom: 12px;
    }
    .title-block h1 { margin: 0; }
    .title-block p  { color: #888; font-size: 9.5pt; margin-top: 6px; }
    @page { margin: 0.8in; }
    """

    import datetime
    full_html = f"""<!DOCTYPE html>
<html><head>
<meta charset="utf-8">
<title>{html_module.escape(title)}</title>
<style>{css}</style>
</head>
<body>
<div class="title-block">
  <h1>{html_module.escape(title)}</h1>
  <p>Generated by Claude Office Assistant &middot; {datetime.datetime.now().strftime("%B %d, %Y")}</p>
</div>
{html_body}
</body></html>"""

    try:
        from weasyprint import HTML as WH

        buf = io.BytesIO()
        WH(string=full_html).write_pdf(buf)
        buf.seek(0)
        return buf
    except ImportError as e:
        logger.warning("WeasyPrint not importable (%s); using ReportLab PDF fallback", e)
        return export_pdf_reportlab(markdown_text, title=title)
    except Exception as e:
        logger.warning("WeasyPrint PDF failed (%s); using ReportLab PDF fallback", e)
        return export_pdf_reportlab(markdown_text, title=title)


# ─── PPTX helpers ────────────────────────────────────────────────────────────

_MD_IMG = re.compile(r"!\[[^\]]*\]\((https?://[^\)]+)\)")
# Allow optional list marker so "- IMAGE: https://..." is stripped (not shown as body text)
_IMG_LINE = re.compile(
    r"^\s*(?:[-*+•]\s+)?(?:IMAGE|IMG|img)\s*:\s*(https?://\S+)\s*$",
    re.I,
)
# Headline lines like "Slide 1: Topic" or "**Slide 1:** **Topic**"; multiline so $ is end-of-line
_SLIDE_HEAD = re.compile(
    r"^\s*\*{0,3}\s*Slide\s+(\d+)\s*:\s*\*{0,3}\s*(.+?)\s*(?:\*{0,3}\s*)?$",
    re.I | re.M,
)
_MD_SLIDE = re.compile(r"^##\s+SLIDE\s+\d+\s*:\s*(.+)$", re.I)


def normalize_slide_markdown_for_pptx(text: str, default_title: str = "Presentation") -> str:
    """
    Convert 'Slide N: headline' prose (common from Claude) into ## SLIDE N: markdown.
    Preserves lines that already use ## SLIDE N: headers.
    """
    t = (text or "").strip()
    if not t:
        return f"# {default_title}\n"
    if re.search(r"(?m)^##\s+SLIDE\s+\d+\s*:", t):
        if not re.search(r"(?m)^#\s+", t):
            return f"# {default_title}\n\n{t}".strip()
        return t

    lines = t.splitlines()
    idx_first = None
    for i, ln in enumerate(lines):
        if _SLIDE_HEAD.match(ln.strip()):
            idx_first = i
            break

    out: List[str] = []
    scan = "\n".join(lines[idx_first:]) if idx_first is not None else t

    if idx_first is not None and idx_first > 0:
        raw_pre = [
            x.strip().lstrip("#").strip()
            for x in lines[:idx_first]
            if x.strip()
        ]
        junk = re.compile(r"^(?:✏️\s*)?edit\s*$|^(?:✦)+\s*$|^claude\s*$|^.{1,6}$", re.I)
        titleish = []
        for p in raw_pre:
            if junk.match(p):
                continue
            if _SLIDE_HEAD.match(p):
                continue
            if re.search(r"slide\s*\d+", p, re.I):
                continue
            if len(p) >= 12:
                titleish.append(p)
        if titleish:
            deck = titleish[-1]
            out.append(f"# {deck}")
            out.append("")
        elif raw_pre:
            cand = raw_pre[-1]
            out.append(f"# {cand if len(cand) >= 6 else default_title}")
            out.append("")
    elif idx_first is None:
        return t  # unknown shape; export_pptx fallback will treat as single block

    spans: List[tuple] = []
    for m in _SLIDE_HEAD.finditer(scan):
        spans.append((int(m.group(1)), m.group(2).strip(), m.start(), m.end()))

    if not spans:
        return t

    for i, (_, headline, _, h_end) in enumerate(spans):
        out.append(f"## SLIDE {spans[i][0]}: {headline}")
        body_start = h_end
        body_end = spans[i + 1][2] if i + 1 < len(spans) else len(scan)
        body = scan[body_start:body_end].strip()
        skip_notes = False
        for raw in body.splitlines():
            s = raw.strip()
            if not s or _SLIDE_HEAD.match(s):
                continue
            ls = s.lower()
            if ls.startswith("speaker notes") or ls.startswith("[notes:"):
                skip_notes = True
            if skip_notes:
                continue
            if s.startswith(("[NOTES:", "[notes:")):
                continue
            if "\t" in s and "|" not in s and not s.startswith("-"):
                parts = [p.strip() for p in s.split("\t") if p.strip()]
                if len(parts) >= 2:
                    s = "- " + " · ".join(parts)
            line_out = s if s.startswith(("-", "*", "•", "+", "|")) else f"- {s}"
            out.append(line_out)
        out.append("")

    return "\n".join(out).strip()


def _extract_slide_image_urls(bullets: list) -> Tuple[List[str], Optional[str]]:
    """Pull first HTTPS image URL from bullets; return (clean_bullets, url or None)."""
    url = None
    cleaned = []
    for b in bullets:
        raw = b.strip()
        ml = _IMG_LINE.match(raw)
        if ml and not url:
            url = ml.group(1).strip().rstrip(").,;")
            continue
        # Bare URL line (common from models): treat as image-only row
        if not url and re.match(r"^\s*(?:[-*+•]\s+)?https?://\S+\s*$", raw):
            url = re.sub(r"^\s*(?:[-*+•]\s+)?", "", raw).strip().rstrip(").,;")
            continue
        m = _MD_IMG.search(raw)
        if m and not url:
            url = m.group(1).strip()
            raw = _MD_IMG.sub("", raw).strip()
            if not raw or raw in "-*+•":
                continue
        if raw:
            cleaned.append(raw if raw.startswith("-") else f"- {raw}")
    return cleaned, url


def _fetch_image_bytes(url: str, timeout: int = 24) -> Optional[bytes]:
    try:
        import requests as req
        if not url.startswith(("http://", "https://")):
            return None
        u = url.split("#")[0].strip()
        host = ""
        try:
            from urllib.parse import urlparse

            host = urlparse(u).hostname or ""
        except Exception:
            pass
        ref = None
        if "unsplash.com" in host:
            ref = "https://unsplash.com/"
        elif host:
            ref = f"https://{host}/"
        headers = {
            "User-Agent": "Mozilla/5.0 (compatible; ClaudeOfficeExport/1.0; +office-assistant)",
            "Accept": "image/avif,image/webp,image/png,image/jpeg,image/*;q=0.8,*/*;q=0.5",
        }
        if ref:
            headers["Referer"] = ref
        r = req.get(u, timeout=timeout, headers=headers, allow_redirects=True)
        if r.ok and not r.content.startswith(b"<htm") and len(r.content) > 100:
            return r.content
    except Exception:
        pass
    return None


def _placeholder_image_bytes(seed: str) -> Optional[bytes]:
    """Lightweight thematic placeholder (no API key)."""
    try:
        import hashlib as hl
        import requests as req
        h = hl.md5(seed.encode("utf-8"), usedforsecurity=False).hexdigest()[:16]
        u = f"https://picsum.photos/seed/{h}/520/340"
        r = req.get(u, timeout=12, headers={"User-Agent": "ClaudeOfficeExport/1.0"})
        if r.ok and len(r.content) > 500:
            return r.content
    except Exception:
        pass
    return None


def _offline_gradient_placeholder_png(seed: str) -> bytes:
    """Tiny two-tone PNG (no network); varies by seed so slides look distinct."""
    import hashlib as hl
    import struct
    import zlib

    d = hl.md5(seed.encode("utf-8"), usedforsecurity=False).digest()
    r1 = 0x14 + d[0] % 50
    g1 = 0x1c + d[1] % 60
    b1 = 0x3a + d[2] % 80
    r2 = min(255, r1 + 40)
    g2 = min(255, g1 + 55)
    b2 = min(255, b1 + 70)
    w, h = 64, 42
    row1 = b"\x00" + bytes([r1, g1, b1]) * w
    row2 = b"\x00" + bytes([r2, g2, b2]) * w
    half = h // 2
    raw_data = row1 * half + row2 * (h - half)

    def chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
    return (
        bytes([137, 80, 78, 71, 13, 10, 26, 10])
        + chunk(b"IHDR", ihdr)
        + chunk(b"IDAT", zlib.compress(raw_data, 9))
        + chunk(b"IEND", b"")
    )


def _pptx_parse_slide_blocks(markdown_text: str, default_title: str) -> List[Tuple[str, List[str], Optional[str]]]:
    """
    Parsed slides: [(title, bullet_lines_with_dash_prefix, optional_image_https_url), ...]
    Optional leading '# Deck Title' adds a dedicated title slide.
    """
    lines = markdown_text.splitlines()
    i = 0
    deck: Optional[str] = None

    while i < len(lines) and not lines[i].strip():
        i += 1
    if i < len(lines) and re.match(r"^#\s+", lines[i].strip()):
        deck = lines[i].lstrip("# ").strip()
        i += 1
    while i < len(lines) and not lines[i].strip():
        i += 1

    hdr_re = re.compile(r"^##\s+SLIDE\s+\d+\s*:\s*(.+)$", re.I)

    blocks: List[Tuple[str, List[str], Optional[str]]] = []

    def flush(acc_title: str, acc_lines: list) -> None:
        bullets_raw = []
        for ln in acc_lines:
            s = ln.strip()
            if not s:
                continue
            low = s.lower()
            if low.startswith("speaker notes") or low.startswith("[notes:"):
                break
            if s.startswith(("[NOTES:", "[notes:")):
                continue
            if s.startswith(("-", "*", "•", "+")):
                bullets_raw.append(s)
            elif s.startswith("|"):
                bullets_raw.append(f"- {s}")
            else:
                bullets_raw.append(f"- {s}")
        bullets, img = _extract_slide_image_urls(bullets_raw)
        blocks.append((acc_title, bullets, img))
    cur_title = ""
    cur_lines: list = []

    while i < len(lines):
        m = hdr_re.match(lines[i].strip())
        if m:
            if cur_title:
                flush(cur_title, cur_lines)
                cur_lines = []
            else:
                cur_lines = []
            cur_title = m.group(1).strip()
        else:
            cur_lines.append(lines[i])
        i += 1
    if cur_title:
        flush(cur_title, cur_lines)

    result: List[Tuple[str, List[str], Optional[str]]] = []
    if deck:
        result.append((deck, [], None))
    for title, bulls, img in blocks:
        if not title:
            title = default_title
        result.append((title, bulls if bulls else [], img))

    if not result and markdown_text.strip():
        chunk = markdown_text.strip()[:2000]
        result.append((default_title, [f"- {chunk}"], None))
    elif not result:
        result.append((default_title, [], None))
    return result


# ─── PPTX Export ─────────────────────────────────────────────────────────────

def export_pptx(markdown_text: str, title: str = "Presentation") -> io.BytesIO:
    _ensure_python_pptx_installed()
    from pptx import Presentation as Prs
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Prs()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Professional deck palette (navy + gold, closer to polished product decks) ─
    DARK_BG   = RGBColor(0x1a, 0x2b, 0x48)
    GOLD      = RGBColor(0xd4, 0xaf, 0x37)
    LIGHT_TXT = RGBColor(0xf5, 0xf6, 0xf8)
    MUTED     = RGBColor(0xb8, 0xc4, 0xdc)
    WHITE     = RGBColor(0xff, 0xff, 0xff)

    TITLE_FACE = "Georgia"
    BODY_FACE  = "Calibri"

    blank_layout = prs.slide_layouts[6]  # Completely blank

    def hex_bg(slide, color: RGBColor):
        """Fill slide background with solid color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_slide_footer(slide, slide_num: int, total_slides: int) -> None:
        """Small page indicator bottom-right (matches common deck UX)."""
        foot = slide.shapes.add_textbox(
            Inches(11.35), Inches(7.05), Inches(1.85), Inches(0.38)
        )
        tf = foot.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"Page {slide_num} / {total_slides}"
        r.font.name = BODY_FACE
        r.font.size = Pt(10)
        r.font.color.rgb = MUTED

    def title_slide(heading: str, subtitle: str = "", slide_num: int = 1, total_slides: int = 1):
        slide = prs.slides.add_slide(blank_layout)
        hex_bg(slide, DARK_BG)

        # Gold bands top + bottom (executive deck style)
        from pptx.util import Emu

        for y in (Inches(0), Inches(7.42)):
            band = slide.shapes.add_shape(1, Inches(0), y, Inches(13.33), Emu(2400))
            band.fill.solid()
            band.fill.fore_color.rgb = GOLD
            band.line.fill.background()

        # Title — centered
        txBox = slide.shapes.add_textbox(Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.35))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = heading.upper() if len(heading) < 80 else heading
        run.font.name  = TITLE_FACE
        run.font.size  = Pt(44)
        run.font.bold  = True
        run.font.color.rgb = WHITE

        if subtitle:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.name  = BODY_FACE
            r2.font.size  = Pt(18)
            r2.font.color.rgb = GOLD

        add_slide_footer(slide, slide_num, total_slides)

    def content_slide(
        heading: str,
        bullets: list,
        image_bytes: Optional[bytes] = None,
        slide_num: int = 2,
        total_slides: int = 2,
    ):
        slide = prs.slides.add_slide(blank_layout)
        hex_bg(slide, DARK_BG)

        # Top gold band
        from pptx.util import Emu

        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Emu(2800))
        bar.fill.solid()
        bar.fill.fore_color.rgb = GOLD
        bar.line.fill.background()

        # Heading
        hbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.42), Inches(12), Inches(0.95))
        hf = hbox.text_frame
        hp = hf.paragraphs[0]
        hr = hp.add_run()
        hr.text = heading
        hr.font.name = TITLE_FACE
        hr.font.size = Pt(28)
        hr.font.bold = True
        hr.font.color.rgb = DARK_BG

        # Gold underline under title
        line = slide.shapes.add_shape(
            1, Inches(0.6), Inches(1.28), Inches(7.2) if image_bytes else Inches(12.1), Emu(2800)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.line.fill.background()

        # Bullets (narrower when image present)
        bullet_w = Inches(6.85) if image_bytes else Inches(12.1)

        bBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), bullet_w, Inches(5.6))
        btf = bBox.text_frame
        btf.word_wrap = True

        for idx, bullet in enumerate(bullets):
            text = bullet.lstrip("-•* ")
            p = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
            p.space_before = Pt(4)
            p.space_after  = Pt(4)
            run = p.add_run()
            # Sub-bullet detection (starts with spaces or double dash)
            if bullet.startswith(("  ", "\t", "    ")):
                run.text = f"   › {text}"
                run.font.size  = Pt(15)
                run.font.color.rgb = MUTED
            else:
                run.text = f"• {text}"
                run.font.size  = Pt(18)
                run.font.color.rgb = LIGHT_TXT
            run.font.name = BODY_FACE

        add_slide_footer(slide, slide_num, total_slides)

        if image_bytes:
            import os
            import tempfile

            ext = ".png" if image_bytes[:4] == b"\x89PNG" else ".jpg"
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(image_bytes)
                path = tmp.name
            try:
                slide.shapes.add_picture(
                    path,
                    Inches(7.15),
                    Inches(1.35),
                    width=Inches(5.65),
                )
            finally:
                try:
                    os.unlink(path)
                except OSError:
                    pass

    # ── Parse markdown into slides (normalize Slide N: prose + ## SLIDE blocks) ─
    md = normalize_slide_markdown_for_pptx(markdown_text, title)
    slides_list = _pptx_parse_slide_blocks(md, title)

    if not slides_list:
        slides_list = [(title, [f"- {markdown_text[:300]}"], None)]

    # Title card: explicit '# deck' row is empty bullets; else use API title
    if slides_list[0][1] == []:
        opening_title = slides_list[0][0]
        content_rows = slides_list[1:]
    else:
        opening_title = title
        content_rows = slides_list

    total_slides = 1 + len(content_rows)
    title_slide(opening_title, "", slide_num=1, total_slides=total_slides)

    for i, (heading, bullets, img_url) in enumerate(content_rows):
        pic: Optional[bytes] = None
        if img_url:
            pic = _fetch_image_bytes(img_url)
        if pic is None:
            pic = _placeholder_image_bytes(f"{heading}|{title}")
        if pic is None:
            pic = _offline_gradient_placeholder_png(f"{heading}|{title}")
        content_slide(
            heading,
            bullets if bullets else [],
            pic,
            slide_num=i + 2,
            total_slides=total_slides,
        )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
